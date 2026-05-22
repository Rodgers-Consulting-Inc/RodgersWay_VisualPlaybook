"""
Export RodgersWay Visual Playbook → PowerPoint (.pptx)

Run from the repo root:
    python tools/export_to_pptx.py
Produces: RodgersWay_Playbook.pptx
"""

import json, os, textwrap
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── paths ────────────────────────────────────────────────────────────────────
ROOT   = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DATA   = os.path.join(ROOT, "data", "library.json")
OUTPUT = os.path.join(ROOT, "RodgersWay_Playbook.pptx")

# ── slide dimensions (widescreen 16:9) ───────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

# ── brand colours ────────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY  = RGBColor(0x2B, 0x2B, 0x2B)
MID_GRAY   = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
LESSON_BG  = RGBColor(0xF7, 0xF9, 0xE8)   # soft green tint
LESSON_FG  = RGBColor(0x3A, 0x4A, 0x00)


def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def solid_fill(shape, colour: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour


def add_textbox(slide, left, top, width, height, text, font_size,
                bold=False, colour=DARK_GRAY, align=PP_ALIGN.LEFT,
                wrap=True, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = colour
    run.font.italic = italic
    return txb


def add_rect(slide, left, top, width, height, fill: RGBColor, line=None):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    solid_fill(shape, fill)
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


# ── cover slide ──────────────────────────────────────────────────────────────
def make_cover(prs, title, subtitle, tagline):
    layout = prs.slide_layouts[6]   # blank
    slide  = prs.slides.add_slide(layout)

    # dark background
    add_rect(slide, 0, 0, W, H, DARK_GRAY)

    # accent bar (Rodgers green)
    add_rect(slide, 0, Inches(4.6), W, Inches(0.08), RGBColor(0xAF, 0xC9, 0x35))

    add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
                title, 40, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.1), Inches(11), Inches(0.7),
                subtitle, 20, colour=RGBColor(0xCC, 0xCC, 0xCC),
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
                tagline, 13, colour=RGBColor(0xAF, 0xC9, 0x35),
                align=PP_ALIGN.CENTER, italic=True)


# ── module divider slide ──────────────────────────────────────────────────────
def make_module_cover(prs, module):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    colour = hex_to_rgb(module.get("color", "#4A90D9"))

    add_rect(slide, 0, 0, W, H, colour)

    # icon / number chip
    icon_txt = module.get("icon", "")
    if icon_txt:
        add_rect(slide, Inches(1), Inches(2.0), Inches(0.9), Inches(0.9),
                 RGBColor(0xFF, 0xFF, 0xFF))
        add_textbox(slide, Inches(1), Inches(2.0), Inches(0.9), Inches(0.9),
                    icon_txt, 22, bold=True,
                    colour=colour, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2.2), Inches(2.0), Inches(9.5), Inches(1.1),
                module["title"], 34, bold=True, colour=WHITE)

    slide_count = len(module["slides"])
    add_textbox(slide, Inches(2.2), Inches(3.2), Inches(9.5), Inches(0.5),
                f"{slide_count} slide{'s' if slide_count != 1 else ''}", 16,
                colour=RGBColor(0xFF, 0xFF, 0xFF))


# ── content slide ─────────────────────────────────────────────────────────────
TITLE_H    = Inches(0.72)
TITLE_TOP  = Inches(0.0)
BAR_H      = Inches(0.06)
MARGIN     = Inches(0.28)
CONTENT_TOP = Inches(0.78)
FOOTER_H   = Inches(0.9)


def make_content_slide(prs, slide_data, module_colour: RGBColor):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)

    # ── title bar ────────────────────────────────────────────────────────────
    add_rect(slide, 0, TITLE_TOP, W, TITLE_H, DARK_GRAY)
    add_rect(slide, 0, TITLE_H,   W, BAR_H,   module_colour)

    title = slide_data.get("title", "")
    add_textbox(slide, MARGIN, Inches(0.08), W - MARGIN * 2, Inches(0.58),
                title, 18, bold=True, colour=WHITE)

    # ── parse blocks ─────────────────────────────────────────────────────────
    images     = []
    key_lesson = ""
    teaching   = {}

    for b in slide_data.get("blocks", []):
        if b["type"] == "imagePair":
            images = [i for i in b["images"] if i.get("src")]
        elif b["type"] == "keyLesson":
            key_lesson = b.get("text", "")
        elif b["type"] == "teaching":
            teaching = b

    # ── layout geometry ───────────────────────────────────────────────────────
    has_lesson = bool(key_lesson.strip())
    footer_top = H - FOOTER_H - Inches(0.15) if has_lesson else None
    img_bottom = footer_top - Inches(0.12) if has_lesson else H - Inches(0.15)
    img_area_h = img_bottom - CONTENT_TOP

    # ── images ────────────────────────────────────────────────────────────────
    if len(images) == 2:
        slot_w = (W - MARGIN * 3) / 2
        for idx, img in enumerate(images):
            left = MARGIN + idx * (slot_w + MARGIN)
            _place_image(slide, img, left, CONTENT_TOP, slot_w, img_area_h)
    elif len(images) == 1:
        img_w = W - MARGIN * 2
        _place_image(slide, images[0], MARGIN, CONTENT_TOP, img_w, img_area_h)
    else:
        # no image — show a placeholder box
        add_rect(slide, MARGIN, CONTENT_TOP, W - MARGIN * 2, img_area_h,
                 LIGHT_GRAY)
        add_textbox(slide, MARGIN, CONTENT_TOP + img_area_h / 2 - Inches(0.25),
                    W - MARGIN * 2, Inches(0.5),
                    "[ No image available for this slide ]",
                    13, colour=MID_GRAY, align=PP_ALIGN.CENTER)

    # ── key lesson footer ─────────────────────────────────────────────────────
    if has_lesson:
        add_rect(slide, MARGIN, footer_top, W - MARGIN * 2, FOOTER_H, LESSON_BG)
        add_textbox(slide, MARGIN + Inches(0.15), footer_top + Inches(0.07),
                    Inches(1.0), Inches(0.38),
                    "KEY LESSON", 8, bold=True, colour=LESSON_FG)
        add_textbox(slide, MARGIN + Inches(0.15),
                    footer_top + Inches(0.32),
                    W - MARGIN * 2 - Inches(0.3), Inches(0.52),
                    key_lesson, 11, colour=LESSON_FG)

    # ── speaker notes ─────────────────────────────────────────────────────────
    notes_parts = []

    narration = slide_data.get("narration", {})
    script    = narration.get("script", "").strip()
    if script:
        notes_parts.append("NARRATION SCRIPT\n" + script)

    t_lines = [teaching.get(f"t{i}", "").strip() for i in range(1, 6)
               if teaching.get(f"t{i}", "").strip()]
    if t_lines:
        notes_parts.append("TEACHING NOTES\n" +
                           "\n".join(f"• {t}" for t in t_lines))

    audio = narration.get("audio", "")
    if audio:
        notes_parts.append(f"AUDIO FILE: {audio}")

    if notes_parts:
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = "\n\n".join(notes_parts)


def _place_image(slide, img_data, left, top, max_w, max_h):
    """Add an image, preserving aspect ratio within the given box."""
    from PIL import Image as PILImage

    src = img_data.get("src", "")
    label = img_data.get("label", "")
    img_type = img_data.get("type", "")

    path = os.path.join(ROOT, src) if src else None

    if path and os.path.exists(path):
        try:
            with PILImage.open(path) as im:
                iw, ih = im.size
            ratio = iw / ih
            if max_w / ratio <= max_h:
                draw_w = max_w
                draw_h = Emu(int(max_w / ratio))
            else:
                draw_h = max_h
                draw_w = Emu(int(max_h * ratio))
            # centre within the slot
            offset_x = (max_w - draw_w) // 2
            offset_y = (max_h - draw_h) // 2
            slide.shapes.add_picture(path,
                                     left + offset_x, top + offset_y,
                                     draw_w, draw_h)
        except Exception as e:
            _image_error_box(slide, left, top, max_w, max_h, str(e))
    else:
        _image_error_box(slide, left, top, max_w, max_h, src or "missing")

    # label + type caption strip
    cap_top = top + max_h + Inches(0.04)
    if img_type:
        add_textbox(slide, left, cap_top, max_w, Inches(0.25),
                    img_type.upper(), 8, bold=True, colour=MID_GRAY)


def _image_error_box(slide, left, top, w, h, msg):
    add_rect(slide, left, top, w, h, LIGHT_GRAY)
    add_textbox(slide, left + Inches(0.1), top + h / 2 - Inches(0.2),
                w - Inches(0.2), Inches(0.4),
                f"[ Image not found: {msg} ]", 10,
                colour=MID_GRAY, align=PP_ALIGN.CENTER)


# ── main ─────────────────────────────────────────────────────────────────────
def main():
    try:
        from PIL import Image  # noqa — just checking
    except ImportError:
        import subprocess, sys
        subprocess.check_call([sys.executable, "-m", "pip", "install",
                               "Pillow", "-q"])

    with open(DATA, encoding="utf-8") as f:
        data = json.load(f)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    make_cover(prs,
               data.get("title", "Visual Playbook"),
               data.get("subtitle", ""),
               data.get("tagline", ""))

    for module in data["modules"]:
        colour = hex_to_rgb(module.get("color", "#4A90D9"))
        make_module_cover(prs, module)
        for slide_data in module["slides"]:
            make_content_slide(prs, slide_data, colour)

    prs.save(OUTPUT)
    total = sum(len(m["slides"]) for m in data["modules"])
    print(f"Saved: {OUTPUT}")
    print(f"Slides: 1 cover + {len(data['modules'])} module covers"
          f" + {total} content slides = "
          f"{1 + len(data['modules']) + total} total")


if __name__ == "__main__":
    main()
